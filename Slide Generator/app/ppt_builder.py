"""Generate a PPTX deck from Datapizza slide data."""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .gemini_assets import GeminiAssetManager
from .slide_models import DatapizzaSlide, SlideDeck, SlideType

try:
    import cairosvg
except Exception:  # pragma: no cover
    cairosvg = None

try:  # pragma: no cover
    from PIL import Image
except Exception:  # pragma: no cover
    Image = None


@dataclass(frozen=True)
class SlideStyle:
    background_hex: str = "#FFFFFF"
    title_blue_hex: str = "#1B64F5"
    brand_red_hex: str = "#C64336"
    text_hex: str = "#111F2D"
    body_hex: str = "#0D1F2E"
    muted_hex: str = "#737373"
    border_hex: str = "#D0D5DD"
    light_hex: str = "#ECEFF2"
    prompt_accent_hex: str = "#FF5C00"
    title_font: str = "Poppins"
    body_font: str = "Inter"

    @staticmethod
    def _hex_to_rgb(value: str) -> RGBColor:
        value = value.lstrip("#")
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))

    def rgb(self, value: str) -> RGBColor:
        return self._hex_to_rgb(value)


class SlideDeckWriter:
    """Create a Datapizza-inspired PPTX document from slide JSON."""

    def __init__(
        self,
        *,
        style: SlideStyle,
        logo_path: Path,
        cover_logo_path: Path | None = None,
        svg_output_dir: Path | None = None,
        gemini_asset_manager: GeminiAssetManager | None = None,
    ):
        self.style = style
        self.logo_path = self._prepare_image_path(logo_path)
        self.cover_logo_path = self._prepare_image_path(cover_logo_path or logo_path)
        self.svg_output_dir = svg_output_dir
        self.gemini_asset_manager = gemini_asset_manager

    def _prepare_image_path(self, path: Path | None) -> Path | None:
        if not path or not path.exists() or Image is None:
            return path
        try:
            with Image.open(path) as img:
                if img.format == "ICO":
                    converted = path.with_name(f"{path.stem}_converted.png")
                    img.save(converted, format="PNG")
                    return converted
        except Exception:
            return path
        return path

    def build(self, deck: SlideDeck, output_path: Path) -> Path:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for index, slide in enumerate(deck.slides, start=1):
            self._add_slide(prs, slide, index)

        output_path = output_path if output_path.suffix else output_path.with_suffix(".pptx")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        return output_path

    def _add_slide(self, prs: Presentation, slide_data: DatapizzaSlide, slide_index: int) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_background(slide)

        if slide_data.slide_type == SlideType.TITLE:
            self._render_title_slide(prs, slide, slide_data)
        elif slide_data.slide_type == SlideType.SECTION:
            self._render_section_slide(prs, slide, slide_data)
        elif slide_data.slide_type == SlideType.CLOSING:
            self._render_closing_slide(prs, slide, slide_data)
        elif slide_data.slide_type == SlideType.BULLETS:
            self._render_bullets_slide(prs, slide, slide_data, slide_index)
        else:
            self._render_content_slide(prs, slide, slide_data, slide_index)

        if slide_data.speaker_notes:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = slide_data.speaker_notes
            notes_frame.paragraphs[0].font.name = self.style.body_font
            notes_frame.paragraphs[0].font.size = Pt(12)

    def _set_background(self, slide) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.style.rgb(self.style.background_hex)

    def _add_logo(self, slide, prs: Presentation, *, bottom: bool = False, icon: bool = False) -> None:
        path = self.cover_logo_path if icon else self.logo_path
        if not path or not path.exists():
            return
        width = Inches(1.35 if not icon else 0.85)
        left = prs.slide_width - width - Inches(0.45)
        top = prs.slide_height - Inches(1.0) if bottom else Inches(0.35)
        slide.shapes.add_picture(str(path), left, top, width=width)

    def _add_title_bar(self, slide) -> None:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.85),
            Inches(1.65),
            Inches(2.4),
            Pt(4),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.style.rgb(self.style.brand_red_hex)
        shape.line.fill.background()

    def _render_title_slide(self, prs: Presentation, slide, slide_data: DatapizzaSlide) -> None:
        self._add_logo(slide, prs)
        self._add_cover_icon(slide, prs)
        self._add_corner_grid(slide)

        label_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.45), Inches(4.5), Inches(0.35))
        label_frame = label_box.text_frame
        label_frame.text = "DATAPIZZA AI4BUILDERS"
        label_run = label_frame.paragraphs[0].runs[0]
        label_run.font.name = self.style.title_font
        label_run.font.size = Pt(10)
        label_run.font.bold = True
        label_run.font.color.rgb = self.style.rgb(self.style.title_blue_hex)

        title_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.85), Inches(8.8), Inches(2.6))
        frame = title_box.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        p.text = slide_data.title
        run = p.runs[0]
        run.font.name = self.style.title_font
        run.font.size = Pt(34)
        run.font.bold = True
        run.font.color.rgb = self.style.rgb(self.style.brand_red_hex)

        if slide_data.subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1.1), Inches(4.4), Inches(7.9), Inches(1.1)
            )
            sub_frame = sub_box.text_frame
            sub_frame.text = slide_data.subtitle
            sub_run = sub_frame.paragraphs[0].runs[0]
            sub_run.font.name = self.style.body_font
            sub_run.font.size = Pt(17)
            sub_run.font.color.rgb = self.style.rgb(self.style.text_hex)

    def _render_section_slide(self, prs: Presentation, slide, slide_data: DatapizzaSlide) -> None:
        self._add_logo(slide, prs)
        self._add_corner_grid(slide)

        label = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(2.2), Inches(0.4))
        label_frame = label.text_frame
        label_frame.text = "SECTION"
        label_run = label_frame.paragraphs[0].runs[0]
        label_run.font.name = self.style.title_font
        label_run.font.size = Pt(10)
        label_run.font.bold = True
        label_run.font.color.rgb = self.style.rgb(self.style.title_blue_hex)

        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.28), Inches(8.4), Inches(1.7))
        title_frame = title_box.text_frame
        title_frame.text = slide_data.title
        title_run = title_frame.paragraphs[0].runs[0]
        title_run.font.name = self.style.title_font
        title_run.font.size = Pt(30)
        title_run.font.bold = True
        title_run.font.color.rgb = self.style.rgb(self.style.brand_red_hex)

        if slide_data.subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.45), Inches(7.6), Inches(1.0))
            sub_frame = sub_box.text_frame
            sub_frame.text = slide_data.subtitle
            sub_run = sub_frame.paragraphs[0].runs[0]
            sub_run.font.name = self.style.body_font
            sub_run.font.size = Pt(15)
            sub_run.font.color.rgb = self.style.rgb(self.style.text_hex)

        arrow_box = slide.shapes.add_textbox(Inches(11.2), Inches(5.2), Inches(1.0), Inches(0.5))
        arrow_frame = arrow_box.text_frame
        arrow_frame.text = ">>"
        arrow_run = arrow_frame.paragraphs[0].runs[0]
        arrow_run.font.name = self.style.title_font
        arrow_run.font.size = Pt(22)
        arrow_run.font.bold = True
        arrow_run.font.color.rgb = self.style.rgb(self.style.title_blue_hex)

    def _render_closing_slide(self, prs: Presentation, slide, slide_data: DatapizzaSlide) -> None:
        self._add_logo(slide, prs)
        self._add_logo(slide, prs, bottom=True, icon=True)

        title_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.35), Inches(5.5), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.text = slide_data.title
        title_run = title_frame.paragraphs[0].runs[0]
        title_run.font.name = self.style.title_font
        title_run.font.size = Pt(30)
        title_run.font.bold = True
        title_run.font.color.rgb = self.style.rgb(self.style.brand_red_hex)

        if slide_data.body:
            body_box = slide.shapes.add_textbox(Inches(1.1), Inches(3.55), Inches(5.8), Inches(1.1))
            body_frame = body_box.text_frame
            body_frame.text = slide_data.body
            body_run = body_frame.paragraphs[0].runs[0]
            body_run.font.name = self.style.body_font
            body_run.font.size = Pt(16)
            body_run.font.color.rgb = self.style.rgb(self.style.body_hex)

        self._render_prompt_chip(slide, "Closing")

    def _render_content_slide(
        self, prs: Presentation, slide, slide_data: DatapizzaSlide, slide_index: int
    ) -> None:
        image_path = self._resolve_visual_path(slide_index)
        if image_path:
            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            return
        self._render_full_slide_placeholder(prs, slide, slide_data)

    def _render_bullets_slide(
        self, prs: Presentation, slide, slide_data: DatapizzaSlide, slide_index: int
    ) -> None:
        self._add_logo(slide, prs)
        self._add_title_and_subtitle(slide, slide_data)

        top = Inches(2.15)
        for idx, bullet in enumerate(slide_data.bullets[:4]):
            row_top = top + Inches(0.9 * idx)
            badge = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(0.9),
                row_top,
                Inches(0.28),
                Inches(0.28),
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = self.style.rgb(self.style.title_blue_hex)
            badge.line.fill.background()

            box = slide.shapes.add_textbox(Inches(1.3), row_top - Inches(0.05), Inches(4.8), Inches(0.45))
            frame = box.text_frame
            frame.text = bullet
            run = frame.paragraphs[0].runs[0]
            run.font.name = self.style.body_font
            run.font.size = Pt(18)
            run.font.color.rgb = self.style.rgb(self.style.body_hex)

        self._render_visual_or_placeholder(slide, slide_data, slide_index, is_bullets=True)

    def _add_cover_icon(self, slide, prs: Presentation) -> None:
        path = self.cover_logo_path
        if not path or not path.exists():
            return
        slide.shapes.add_picture(
            str(path),
            prs.slide_width - Inches(3.2),
            Inches(3.75),
            width=Inches(1.95),
        )

    def _add_title_and_subtitle(self, slide, slide_data: DatapizzaSlide) -> None:
        title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.75), Inches(8.6), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = slide_data.title
        title_run = title_frame.paragraphs[0].runs[0]
        title_run.font.name = self.style.title_font
        title_run.font.size = Pt(24)
        title_run.font.bold = True
        title_run.font.color.rgb = self.style.rgb(self.style.title_blue_hex)
        self._add_title_bar(slide)

        if slide_data.subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.8), Inches(6.5), Inches(0.45))
            sub_frame = sub_box.text_frame
            sub_frame.text = slide_data.subtitle
            sub_run = sub_frame.paragraphs[0].runs[0]
            sub_run.font.name = self.style.body_font
            sub_run.font.size = Pt(11)
            sub_run.font.color.rgb = self.style.rgb(self.style.muted_hex)

    def _render_visual_or_placeholder(
        self,
        slide,
        slide_data: DatapizzaSlide,
        slide_index: int,
        *,
        is_bullets: bool,
    ) -> None:
        left = Inches(7.0)
        top = Inches(2.0)
        width = Inches(5.25)
        height = Inches(3.3 if is_bullets else 3.6)

        image_path = self._resolve_visual_path(slide_index)
        if image_path:
            slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
            return

        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            top,
            width,
            height,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.style.rgb(self.style.light_hex)
        card.line.color.rgb = self.style.rgb(self.style.border_hex)

        prompt_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.35), width - Inches(0.6), height - Inches(0.5))
        prompt_frame = prompt_box.text_frame
        prompt_frame.clear()
        prompt_frame.word_wrap = True
        prompt_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

        header = prompt_frame.paragraphs[0]
        header.text = "Gemini SVG atteso"
        header.font.name = self.style.title_font
        header.font.size = Pt(12)
        header.font.bold = True
        header.font.color.rgb = self.style.rgb(self.style.prompt_accent_hex)

        body = prompt_frame.add_paragraph()
        body.text = slide_data.image_prompt
        body.font.name = self.style.body_font
        body.font.size = Pt(10)
        body.font.color.rgb = self.style.rgb(self.style.body_hex)

    def _render_full_slide_placeholder(self, prs: Presentation, slide, slide_data: DatapizzaSlide) -> None:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.35),
            Inches(0.35),
            prs.slide_width - Inches(0.7),
            prs.slide_height - Inches(0.7),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.style.rgb(self.style.light_hex)
        card.line.color.rgb = self.style.rgb(self.style.border_hex)

        title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.9), Inches(11.0), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = slide_data.title
        title_run = title_frame.paragraphs[0].runs[0]
        title_run.font.name = self.style.title_font
        title_run.font.size = Pt(26)
        title_run.font.bold = True
        title_run.font.color.rgb = self.style.rgb(self.style.title_blue_hex)

        if slide_data.body:
            body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.85), Inches(4.8), Inches(1.4))
            body_frame = body_box.text_frame
            body_frame.word_wrap = True
            body_frame.text = slide_data.body
            body_run = body_frame.paragraphs[0].runs[0]
            body_run.font.name = self.style.body_font
            body_run.font.size = Pt(17)
            body_run.font.color.rgb = self.style.rgb(self.style.body_hex)

        prompt_box = slide.shapes.add_textbox(Inches(5.9), Inches(1.85), Inches(5.4), Inches(3.8))
        prompt_frame = prompt_box.text_frame
        prompt_frame.word_wrap = True
        prompt_frame.text = slide_data.image_prompt
        prompt_run = prompt_frame.paragraphs[0].runs[0]
        prompt_run.font.name = self.style.body_font
        prompt_run.font.size = Pt(11)
        prompt_run.font.color.rgb = self.style.rgb(self.style.body_hex)

    def _render_prompt_chip(self, slide, text: str) -> None:
        chip = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.95),
            Inches(5.9),
            Inches(1.2),
            Inches(0.36),
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = self.style.rgb(self.style.light_hex)
        chip.line.fill.background()
        frame = chip.text_frame
        frame.text = text
        run = frame.paragraphs[0].runs[0]
        run.font.name = self.style.title_font
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = self.style.rgb(self.style.title_blue_hex)
        frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_corner_grid(self, slide) -> None:
        for row in range(4):
            for col in range(4):
                dot = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.OVAL,
                    Inches(10.4 + col * 0.22),
                    Inches(1.2 + row * 0.22),
                    Inches(0.06),
                    Inches(0.06),
                )
                dot.fill.solid()
                dot.fill.fore_color.rgb = self.style.rgb(self.style.border_hex)
                dot.line.fill.background()

    def _resolve_visual_path(self, slide_index: int) -> Path | None:
        if self.gemini_asset_manager:
            svg_path = self.gemini_asset_manager.resolve_svg_for_slide(slide_index)
            png_path = self._materialize_svg(svg_path, slide_index)
            if png_path:
                return png_path
        return None

    def _materialize_svg(self, svg_path: Path | None, slide_index: int) -> Path | None:
        if not svg_path or not svg_path.exists() or not self.svg_output_dir or not cairosvg:
            return None

        self.svg_output_dir.mkdir(parents=True, exist_ok=True)
        png_path = self.svg_output_dir / f"slide_{slide_index:02d}.png"
        try:
            cairosvg.svg2png(url=str(svg_path), write_to=str(png_path))
        except Exception:
            return None
        return png_path if png_path.exists() else None
